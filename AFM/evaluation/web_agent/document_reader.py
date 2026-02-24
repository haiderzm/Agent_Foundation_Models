import os
import json
import csv
import docx
import pptx
import PyPDF2
import pandas as pd
import openpyxl
from typing import Union, List, Dict, Any, Optional


class DocumentReader:
    """
    A GAIA-compatible document reader for:
    TXT, PDF, DOCX, PPTX, CSV, XLSX, JSON, JSON-LD
    """

    def __init__(self):
        pass  # No file path at init

    def read(self, file_path: str) -> Union[List[str], List[Dict], Dict]:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"{file_path} does not exist.")
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".txt" or ext == ".py":
            return self._read_txt(file_path)
        elif ext == ".pdf":
            return self._read_pdf(file_path)
        elif ext == ".docx":
            return self._read_docx(file_path)
        elif ext == ".pptx":
            return self._read_pptx(file_path)
        elif ext in [".csv", ".xlsx"]:
            return self._read_table(file_path)
        elif ext in [".json", ".jsonld"]:
            return self._read_json(file_path)
        elif ext == ".pdb":
            return self._read_pdb(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    # ----------------------------
    # TXT
    # ----------------------------
    def _read_txt(self, file_path: str) -> List[str]:
        with open(file_path, "r", encoding="utf-8") as f:
            return [line.rstrip("\n") for line in f]

    # ----------------------------
    # PDF
    # ----------------------------
    def _read_pdf(self, file_path: str) -> List[str]:
        pages = []
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                pages.append(text if text else "")
        return pages

    # ----------------------------
    # DOCX
    # ----------------------------
    # def _read_docx(self, file_path: str) -> List[str]:
    #     doc = docx.Document(file_path)
    #     return [p.text for p in doc.paragraphs if p.text.strip()]

    def _read_docx(self, file_path: str) -> Dict[str, Any]:
        """
        Returns only data present in the DOCX:
        - paragraphs: non-empty paragraph texts (in order they appear in doc.paragraphs)
        - tables: list of tables; each table is a list of rows; each row is a list of cell texts
        """
        doc = docx.Document(file_path)

        paragraphs: List[str] = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]

        tables: List[List[List[str]]] = []
        for table in doc.tables:
            trows: List[List[str]] = []
            for row in table.rows:
                # cell.text may include internal newlines; normalize whitespace but don't invent structure
                cells = [" ".join(cell.text.split()).strip() for cell in row.cells]
                # keep even if some cells are empty, but drop fully-empty rows
                if any(cells):
                    trows.append(cells)
            if trows:
                tables.append(trows)

        return {
            "type": "docx",
            "path": file_path,
            "paragraphs": paragraphs,
            "tables": tables,
        }

    # ----------------------------
    # PPTX
    # ----------------------------
    def _read_pptx(self, file_path: str) -> List[str]:
        prs = pptx.Presentation(file_path)
        slides_text = []
        for slide in prs.slides:
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())
            slides_text.append("\n".join(text))
        return slides_text

    # ----------------------------
    # CSV / XLSX
    # ----------------------------
    def _read_table(self, file_path: str):
        if file_path.lower().endswith(".csv"):
            with open(file_path, newline="", encoding="utf-8") as f:
                reader = csv.reader(f)
                rows: List[List[str]] = [row for row in reader]

            header = rows[0] if rows else []
            data = rows[1:] if len(rows) > 1 else []

            return {
                "type": "csv",
                "path": file_path,
                "header": header,
                "shape": {"rows": len(data), "cols": len(header)},
                "rows": rows,  # includes header as first row (matches what you printed)
            }
        else:  # XLSX
            return self._read_xlsx(file_path)
        

    def _read_xlsx(self, file_path: str) -> Dict[str, Any]:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        out: Dict[str, Any] = {}

        for ws in wb.worksheets:
            max_row, max_col = ws.max_row, ws.max_column
            grid = []

            for r in range(1, max_row + 1):
                row_cells = []
                for c in range(1, max_col + 1):
                    cell = ws.cell(r, c)

                    # Normalize value (avoid 1.0, 41671.0, etc.)
                    val = cell.value
                    if isinstance(val, float) and val.is_integer():
                        val = int(val)

                    # Extract fill color (RGB only), but only if it's a real solid fill.
                    # This avoids treating "no fill" defaults as black (000000).
                    fill_rgb = None
                    try:
                        fill = cell.fill
                        if fill and getattr(fill, "patternType", None) == "solid":
                            for color in (fill.fgColor, fill.bgColor):
                                if color and color.type == "rgb" and color.rgb:
                                    rgb = color.rgb.upper()
                                    rgb6 = rgb[2:] if len(rgb) == 8 else rgb  # ARGB -> RGB

                                    # Filter common "default/no-fill" blacks
                                    if rgb in ("00000000", "FF000000") or rgb6 == "000000":
                                        continue

                                    fill_rgb = rgb6
                                    break
                    except Exception:
                        pass

                    row_cells.append({
                        "coord": cell.coordinate,
                        "row": r,
                        "col": c,
                        "value": val,
                        "fill": fill_rgb,  # e.g. "F478A7" or None
                    })
                grid.append(row_cells)

            out[ws.title] = {
                "max_row": max_row,
                "max_col": max_col,
                "cells": grid,  # full grid, includes empty cells
            }

        return out  


    # ----------------------------
    # JSON / JSON-LD
    # ----------------------------
    def _read_json(self, file_path: str) -> Dict:
        with open(file_path, "r", encoding="utf-8") as f:
            return json.load(f)
        
def main():
    reader = DocumentReader()
    # fp = "/share/softwares/haider/Agent_Foundation_Models/AFM/data/web_agent/gaia/files/389793a7-ca17-4e82-81cb-2b3a2391b4b9.txt"
    # fp = "/share/softwares/haider/Agent_Foundation_Models/AFM/data/web_agent/gaia/files/a3fbeb63-0e8c-4a11-bff6-0e3b484c3e9c.pptx"
    # fp = "/share/softwares/haider/Agent_Foundation_Models/AFM/data/web_agent/gaia/files/bec74516-02fc-48dc-b202-55e78d0e17cf.jsonld"
    # fp = "/share/softwares/haider/Agent_Foundation_Models/AFM/data/web_agent/gaia/files/67e8878b-5cef-4375-804e-e6291fdbe78a.pdf"
    # fp = "/share/softwares/haider/Agent_Foundation_Models/AFM/data/web_agent/gaia/files/65afbc8a-89ca-4ad5-8d62-355bb401f61d.xlsx"
    # fp = "/share/softwares/haider/Agent_Foundation_Models/AFM/data/web_agent/gaia/files/4d51c4bf-4b0e-4f3d-897b-3f6687a7d9f2.xlsx"
    # fp = "/share/softwares/haider/Agent_Foundation_Models/AFM/data/web_agent/gaia/files/8d46b8d6-b38a-47ff-ac74-cda14cf2d19b.csv"
    # fp = "/share/softwares/haider/Agent_Foundation_Models/AFM/data/web_agent/gaia/files/3da89939-209c-4086-8520-7eb734e6b4ef.xlsx"
    # fp = "/share/softwares/haider/Agent_Foundation_Models/AFM/data/web_agent/gaia/files/67e8878b-5cef-4375-804e-e6291fdbe78a.pdf"
    fp = "/share/softwares/haider/Agent_Foundation_Models/AFM/data/web_agent/gaia/files/cffe0e32-c9a6-4c52-9877-78ceb4aaa9fb.docx"

    content = reader.read(fp)
    print(content)


if __name__ == '__main__':
    main()